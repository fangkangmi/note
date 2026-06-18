#!/usr/bin/env python3
"""Build the LBG 'Agent Skills for Code Review' workshop deck as a .pptx.
Theme mirrors the reveal.js HTML deck (LBG green). Re-skin into the company
template afterwards; speaker notes are on every slide."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ---- palette ----
GREEN  = RGBColor(0x00,0x86,0x5A); DARK  = RGBColor(0x00,0x38,0x2A)
DARK2  = RGBColor(0x00,0x51,0x3C); INK   = RGBColor(0x16,0x24,0x1F)
SOFT   = RGBColor(0xF4,0xF8,0xF6); LINE  = RGBColor(0xDF,0xE7,0xE3)
WHITE  = RGBColor(0xFF,0xFF,0xFF); GLITE = RGBColor(0x7F,0xE3,0xBF)
RED    = RGBColor(0xC0,0x39,0x2A); GREY  = RGBColor(0x9F,0xB4,0xAB)
MUTED  = RGBColor(0x5D,0x6F,0x68)
FONT   = "Segoe UI"; MONO = "Consolas"
MX = Inches(0.7)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()

def tb(s, l, t, w, h):
    box = s.shapes.add_textbox(l, t, w, h); box.text_frame.word_wrap = True
    return box.text_frame

def put(tf, items, dsize=18):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = it.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(it.get('sa', 6)); p.space_before = Pt(it.get('sb', 0))
        r = p.add_run(); r.text = ('•  ' + it['t']) if it.get('bullet') else it['t']
        f = r.font; f.name = FONT; f.size = Pt(it.get('size', dsize))
        f.bold = it.get('bold', False); f.color.rgb = it.get('color', INK)

def header(s, kicker, title, dark=False):
    put(tb(s, MX, Inches(0.5), Inches(11.9), Inches(0.4)),
        [{'t': kicker.upper(), 'size': 12, 'bold': True, 'color': GLITE if dark else GREEN}])
    put(tb(s, MX, Inches(0.95), Inches(11.9), Inches(1.1)),
        [{'t': title, 'size': 33, 'bold': True, 'color': WHITE if dark else DARK}])

def card(s, l, t, w, h, title, body, fill=SOFT, tcolor=GREEN):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.16)
    put(tf, [{'t': title, 'size': 16, 'bold': True, 'color': tcolor, 'sa': 5},
             {'t': body, 'size': 13, 'color': INK}])

# 1 TITLE
s = slide(DARK)
put(tb(s, MX, Inches(1.4), Inches(11.9), Inches(0.4)),
    [{'t': 'LBG ENGINEERING · WORKSHOP', 'size': 13, 'bold': True, 'color': GLITE}])
put(tb(s, MX, Inches(1.9), Inches(11.9), Inches(2.2)),
    [{'t': 'Agent Skills for Code Review', 'size': 52, 'bold': True, 'color': WHITE}])
put(tb(s, MX, Inches(4.0), Inches(11.9), Inches(0.8)),
    [{'t': 'Auto-checking OpenAPI PRs against the LBG APIE standard', 'size': 22, 'bold': True, 'color': GLITE}])
put(tb(s, MX, Inches(4.9), Inches(11.9), Inches(0.6)),
    [{'t': '90 minutes  ·  [date]  ·  [facilitator]', 'size': 16, 'color': GREY}])
notes(s, "Welcome. Goal: decide together whether skill-based review is worth adopting for APIE compliance. "
         "By the end they'll have seen it beat a manual review on the exact task they do weekly. Laptops closed until the hands-on.")

# 2 PROBLEM
s = slide()
header(s, 'The problem', 'Reviewing a spec against APIE is manual toil')
put(tb(s, MX, Inches(2.2), Inches(11.9), Inches(2.6)), [
    {'t': 'Open the PR. Open the APIE standard.', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Scroll. Compare clause by clause. By hand.', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': "Hope you didn't miss one.", 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Every reviewer does it slightly differently.', 'bullet': True, 'size': 20, 'sa': 10}])
put(tb(s, MX, Inches(5.2), Inches(11.9), Inches(0.8)),
    [{'t': 'What if the standard checked itself?', 'size': 26, 'bold': True, 'color': GREEN}])
notes(s, "Ask: 'How long does a proper APIE review of an OpenAPI PR take you?' Then: 'Hands up if you catch every violation, every time.' Few hands — that's the hook.")

# 3 AGENDA
s = slide()
header(s, 'Today · 90 minutes', 'How this runs')
rows = [("12'", "You review it first — find the APIE violations by hand (timed)"),
        ("6'", "Debrief — what did we find, how long did it take"),
        ("10'", "AI review WITHOUT the skill"),
        ("15'", "AI review WITH the APIE skill  — the moment"),
        ("20'", "Your turn — hands-on, in groups"),
        ("15'", "Where this goes — @lbg-reviewer & discussion"),
        ("7'", "Wrap & next steps")]
items = []
for t, d in rows:
    items.append({'t': f"{t}    {d}", 'size': 18, 'sa': 9})
put(tb(s, MX, Inches(2.1), Inches(11.9), Inches(4.6)), items)
notes(s, "Protect the contrast (human vs no-skill vs skill) — it's the payload.")

# 4 SETUP
s = slide()
header(s, 'What we already have', 'Two skills, living in .github/')
card(s, MX, Inches(2.2), Inches(5.75), Inches(1.7), 'APIE standard skill',
     'The LBG API standard, encoded once — the rules an agent checks a spec against.')
card(s, Inches(6.85), Inches(2.2), Inches(5.75), Inches(1.7), 'code-review skill',
     'Detects what changed in the PR and routes the right files to the right standard.')
put(tb(s, MX, Inches(4.3), Inches(11.9), Inches(0.5)),
    [{'t': 'The chain:', 'size': 16, 'bold': True, 'color': INK}])
put(tb(s, MX, Inches(4.8), Inches(11.9), Inches(0.8)),
    [{'t': 'OpenAPI PR diff  →  code-review skill  →  APIE standard skill  →  findings + clause citations',
      'size': 18, 'bold': True, 'color': GREEN}])
notes(s, "Both skills are plain files in .github/, already wired into Copilot. code-review = detection + routing; APIE = the standard.")

# 5 EXERCISE
s = slide(SOFT)
header(s, 'Exercise · 10 min on the clock', "You're the reviewer")
put(tb(s, MX, Inches(2.2), Inches(11.9), Inches(2.4)), [
    {'t': 'Open the PR:  [PR link]', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Open the standard:  [APIE doc link]', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Find every APIE violation. Write each down:  [capture link]', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'No AI yet.', 'bullet': True, 'size': 20, 'sa': 10}])
put(tb(s, MX, Inches(5.1), Inches(11.9), Inches(1.0)),
    [{'t': '10:00', 'size': 54, 'bold': True, 'color': GREEN}])
notes(s, "Start the timer. Hand out the three links. Walk the room — note who reaches for the standard doc vs reviews from memory. This toil is the baseline we beat.")

# 6 THE PR
s = slide()
header(s, 'The change under review', 'Looks small. Is it compliant?')
code = ["   paths:", "+    /Accounts:", "+      post:", "+        summary: Create account",
        "+        responses:", "+          '200':", "+            description: OK",
        "+            content:", "+              application/json:", "+                schema:",
        "+                  type: object", "+                  properties:",
        "+                    AccountID: { type: string }"]
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, Inches(2.1), Inches(9.4), Inches(4.4))
box.fill.solid(); box.fill.fore_color.rgb = SOFT; box.line.color.rgb = LINE; box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
for i, ln in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1); r = p.add_run(); r.text = ln
    r.font.name = MONO; r.font.size = Pt(15)
    r.font.color.rgb = GREEN if ln.startswith('+') else INK
put(tb(s, MX, Inches(6.55), Inches(11.9), Inches(0.5)),
    [{'t': '[swap in your real diff — keep the subtle APIE violations]', 'size': 13, 'color': MUTED}])
notes(s, "Don't reveal the violations yet. Hidden examples: path not lowercase/kebab, no operationId, create should be 201, no 4xx error using the standard Error schema, inline schema vs $ref, property casing. Replace with your ACTUAL APIE rules so a no-skill review can't catch it from general REST knowledge.")

# 7 DEBRIEF + CHART
s = slide()
header(s, 'Debrief', 'What did the room find?')
put(tb(s, MX, Inches(2.3), Inches(5.2), Inches(3.0)), [
    {'t': 'Avg time spent:  [__ min]', 'bullet': True, 'size': 18, 'sa': 10},
    {'t': 'Found the headline violation:  [__ %]', 'bullet': True, 'size': 18, 'sa': 10},
    {'t': 'Most-missed:  [____]', 'bullet': True, 'size': 18, 'sa': 10},
    {'t': 'The gap is the point.', 'size': 22, 'bold': True, 'color': GREEN, 'sb': 14}])
cd = CategoryChartData()
cd.categories = ['Present in PR', 'Found by humans', 'Found by skill']
cd.add_series('Violations', (6, 3, 6))   # <-- EDIT these to your real numbers
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.4), Inches(2.2),
                        Inches(6.2), Inches(4.2), cd)
ch = gf.chart; ch.has_legend = False; ch.has_title = False
pts = ch.plots[0].series[0].points
for idx, col in enumerate((GREY, RED, GREEN)):
    pts[idx].format.fill.solid(); pts[idx].format.fill.fore_color.rgb = col
notes(s, "Fill the numbers live. The bar gap — violations present vs found by humans vs found by skill — is your whole argument. Edit the (6, 3, 6) tuple to match the real PR + room.")

# 8 BET
s = slide(DARK)
header(s, 'Place your bet', 'Will AI with NO skill catch the APIE violation?', dark=True)
put(tb(s, MX, Inches(2.8), Inches(11.9), Inches(1.2)),
    [{'t': 'Yes   /   No', 'size': 40, 'bold': True, 'color': GLITE}])
put(tb(s, MX, Inches(4.1), Inches(11.9), Inches(0.6)),
    [{'t': 'Hands up. Remember your vote.', 'size': 18, 'color': GREY}])
notes(s, "Quick show of hands; count roughly. Sets up the reveal.")

# 9 NO SKILL
s = slide()
header(s, 'AI · without the skill', 'Generic review')
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, Inches(2.1), Inches(11.9), Inches(1.6))
sh.fill.solid(); sh.fill.fore_color.rgb = SOFT; sh.line.color.rgb = GREEN; sh.shadow.inherit = False
put(sh.text_frame, [{'t': '[ live: run Copilot review with NO skill — or paste screenshot ]', 'size': 15, 'color': MUTED, 'align': PP_ALIGN.CENTER}])
put(tb(s, MX, Inches(3.9), Inches(11.9), Inches(1.8)), [
    {'t': 'Comments on style / naming — generically.', 'bullet': True, 'size': 19, 'sa': 8},
    {'t': "Doesn't know the APIE standard. No clause. Misses the org-specific rule.", 'bullet': True, 'size': 19, 'sa': 8}])
put(tb(s, MX, Inches(5.9), Inches(11.9), Inches(0.7)),
    [{'t': "AI alone doesn't know your standard.", 'size': 24, 'bold': True, 'color': GREEN}])
notes(s, "Run live or play recording. Land the line: capable model, but it has never read APIE. (You confirmed it really misses — no domain knowledge.)")

# 10 WITH SKILL
s = slide(DARK2)
header(s, 'AI · with the APIE skill', 'It checks itself', dark=True)
put(tb(s, MX, Inches(2.0), Inches(11.9), Inches(0.6)),
    [{'t': 'PR diff  →  code-review  →  APIE skill  →  violation + clause + fix', 'size': 18, 'bold': True, 'color': GLITE}])
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, Inches(2.8), Inches(11.9), Inches(1.7))
sh.fill.solid(); sh.fill.fore_color.rgb = DARK; sh.line.color.rgb = GLITE; sh.shadow.inherit = False
put(sh.text_frame, [{'t': '[ live: review citing APIE §x.x with suggested fix ]', 'size': 15, 'color': GLITE, 'align': PP_ALIGN.CENTER}])
put(tb(s, MX, Inches(4.8), Inches(11.9), Inches(0.8)),
    [{'t': 'Seconds. Every violation. Cites the exact clause.', 'size': 24, 'bold': True, 'color': GLITE}])
notes(s, "THE MOMENT — slow down. Point at the citation. Re-ask the bet. Caveat: clause citation is usually present but model-dependent — if it doesn't print the clause id this run, narrate it from the standard.")

# 11 SCORECARD
s = slide()
header(s, 'Side by side', 'The contrast')
data = [["", "Time", "Caught APIE violation", "Cited the clause", "Same every time"],
        ["Human (manual)", "~10 min", "some did", "from memory", "varies"],
        ["AI · no skill", "seconds", "No", "No", "generic"],
        ["AI · APIE skill", "seconds", "Yes", "Yes — clause + fix", "Yes — identical"]]
tbl = s.shapes.add_table(4, 5, MX, Inches(2.3), Inches(11.9), Inches(3.4)).table
for c in range(5):
    tbl.columns[c].width = Inches(2.7 if c == 0 else 2.3)
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val if val else " "
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.name = FONT; run.font.size = Pt(13)
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            run.font.color.rgb = WHITE; run.font.bold = True
        elif ri == 3:
            cell.fill.solid(); cell.fill.fore_color.rgb = SOFT
            run.font.color.rgb = DARK; run.font.bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            run.font.color.rgb = INK
notes(s, "This table is the buy-in. Leave it up while you talk through why it matters.")

# 11b BIG NUMBER
s = slide(DARK2)
header(s, 'The difference', 'Manual review → skill review', dark=True)
put(tb(s, MX, Inches(2.6), Inches(11.9), Inches(1.6)),
    [{'t': '~10 min  →  ~10 sec', 'size': 60, 'bold': True, 'color': WHITE}])
put(tb(s, MX, Inches(4.5), Inches(11.9), Inches(0.8)),
    [{'t': 'Same standard. Every PR. With the clause cited.', 'size': 24, 'bold': True, 'color': GLITE}])
notes(s, "A breath. The headline number. Let it sit before 'why it matters'.")

# 12 WHY LBG
s = slide()
header(s, 'Why this matters at LBG', 'Not "cool AI" — governance')
cards = [('Auditability', 'Cites the exact APIE clause → traceable, defensible review.'),
         ('Consistency', 'Same standard applied identically, every team, every PR.'),
         ('Reviewer load', 'Humans focus on logic & risk, not clause-matching.'),
         ('Speed', 'Seconds, not minutes — faster PRs, less back-and-forth.')]
xs = [MX, Inches(6.85)]; ys = [Inches(2.2), Inches(4.15)]
for i, (t, b) in enumerate(cards):
    card(s, xs[i % 2], ys[i // 2], Inches(5.75), Inches(1.7), t, b)
notes(s, "This is the language leadership buys. Auditability first — it's a bank.")

# 13 TWO LENSES (validated)
s = slide(SOFT)
header(s, 'The real finding · pass@k validated', 'Run it with AND without the skill')
card(s, MX, Inches(2.3), Inches(5.75), Inches(2.0), 'With the skill',
     'Deep on APIE — cites clauses, catches standard violations a plain model cannot know. '
     'But it narrows the model’s view to the standard.')
card(s, Inches(6.85), Inches(2.3), Inches(5.75), Inches(2.0), 'Without the skill',
     'Broader, unfocused — catches general issues outside APIE that the skill’s narrow lens skips over.')
put(tb(s, MX, Inches(4.7), Inches(11.9), Inches(0.8)),
    [{'t': 'Both outputs are valuable. Use both lenses.', 'size': 24, 'bold': True, 'color': GREEN}])
notes(s, "The honest, validated core message. We ran pass@k: the skill reliably ADDS APIE depth AND tunnels the model's attention to the standard, missing some general issues the no-skill run catches. That's WHY we show both today: skill = depth on your standard, no-skill = breadth. Keep a human approver over both.")

# 14 HANDS-ON
s = slide(DARK)
header(s, 'Your turn · 20 min', 'Run it yourselves', dark=True)
put(tb(s, MX, Inches(2.3), Inches(11.9), Inches(3.0)), [
    {'t': 'Groups of [4-5]. Sandbox PR: [link]', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10},
    {'t': 'Trigger the review (with the skill). Watch it cite APIE.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10},
    {'t': 'Try to fool it — add a violation, see if it catches it.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10},
    {'t': 'Note: one thing it nailed, one thing it missed.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10}])
notes(s, "Facilitators roam. Collect the 'missed' notes — they feed the skill and the rollout pitch. (Expect some over-flagging; that's the 'both lenses' message in action.)")

# 15 REMOTE AGENT
s = slide(DARK2)
header(s, 'Where this goes', 'Meet @lbg-reviewer', dark=True)
put(tb(s, MX, Inches(2.2), Inches(11.9), Inches(2.4)), [
    {'t': 'Today: on-demand in Copilot, at your desk.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10},
    {'t': 'Now (PoC): comment @lbg-reviewer on any GitHub PR → it runs the APIE review inline.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10},
    {'t': 'Next: automatic on every spec PR. Conversations underway at management level.', 'bullet': True, 'size': 20, 'color': WHITE, 'sa': 10}])
put(tb(s, MX, Inches(5.1), Inches(11.9), Inches(0.7)),
    [{'t': '▶ Live demo', 'size': 24, 'bold': True, 'color': GLITE}])
notes(s, "Do the LIVE demo here — comment @lbg-reviewer on a real PR and let it post the APIE review inline. Be honest it's a PoC; the point is direction: same skills, now a teammate on every PR.")

# 16 DISCUSSION
s = slide()
header(s, "Let's decide together", 'Open questions for @lbg-reviewer')
put(tb(s, MX, Inches(2.2), Inches(11.9), Inches(3.2)), [
    {'t': 'Advisory comment, or block the merge?', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Who signs off for audit — the bot, or a human approver?', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'How do we keep the APIE skill in sync as the standard evolves?', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Which repos / which standards next?', 'bullet': True, 'size': 20, 'sa': 10}])
notes(s, "Capture answers — genuinely useful input for the management track. Assign an owner to each.")

# 17 CTA
s = slide()
header(s, 'Next steps', 'This week')
put(tb(s, MX, Inches(2.3), Inches(11.9), Inches(2.6)), [
    {'t': 'Run the skill on your next OpenAPI PR.', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': 'Send us anything it missed or over-flagged → it makes the skill better.', 'bullet': True, 'size': 20, 'sa': 10},
    {'t': '[owner] is collating feedback for the @lbg-reviewer rollout.', 'bullet': True, 'size': 20, 'sa': 10}])
notes(s, "Concrete, small, this-week asks. End with the owner + where to send feedback.")

# 18 THANKS
s = slide(DARK)
put(tb(s, MX, Inches(2.6), Inches(11.9), Inches(1.5)),
    [{'t': 'Thank you', 'size': 54, 'bold': True, 'color': WHITE}])
put(tb(s, MX, Inches(4.1), Inches(11.9), Inches(0.7)),
    [{'t': 'Questions?', 'size': 26, 'bold': True, 'color': GLITE}])
put(tb(s, MX, Inches(4.9), Inches(11.9), Inches(0.6)),
    [{'t': '[facilitator] · [contact] · skills repo: [.github link]', 'size': 15, 'color': GREY}])
notes(s, "Leave the scorecard or 'two lenses' slide up if Q&A runs long.")

out = "/home/fang/workspace/tmp/skill-demo/lbg-workshop-deck.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
